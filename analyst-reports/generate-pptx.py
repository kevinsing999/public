#!/usr/bin/env python3
"""Generate PowerPoint slide pack from SBC Deployment Options markdown."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import datetime

# --- Colours ---
DARK_BG = RGBColor(0x1B, 0x2A, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
MID_GREY = RGBColor(0xBF, 0xBF, 0xBF)
DARK_TEXT = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_GREEN = RGBColor(0x2D, 0x8B, 0x57)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TABLE_HEADER_BG = RGBColor(0x2E, 0x75, 0xB6)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xFA)
TABLE_WHITE_BG = RGBColor(0xFF, 0xFF, 0xFF)

# --- Glossary per slide ---
GLOSSARY_COMMON = """GLOSSARY OF TERMS:
- SBC: Session Border Controller - network device managing VoIP sessions at network boundaries
- HA: High Availability - architecture eliminating single points of failure
- VIP: Virtual IP - floating IP address that moves between HA pair members during failover
- EIP: Elastic IP - AWS static public IP address
- IAM: Identity and Access Management - AWS permission and access control framework
- VPC: Virtual Private Cloud - isolated virtual network in AWS
- ENI: Elastic Network Interface - virtual network card attached to an EC2 instance
- AZ: Availability Zone - physically separate data centre within an AWS region
- VRRP: Virtual Router Redundancy Protocol - standard protocol for automatic IP failover
- RTO: Recovery Time Objective - maximum acceptable downtime after a failure
- MSP: Managed Service Provider - third-party IT operations partner
- TTO: Technology Team Operations - internal cloud platform team
- SIP: Session Initiation Protocol - signalling protocol for voice/video sessions
- DNS TTL: Domain Name System Time To Live - cache duration for DNS records
- PSTN: Public Switched Telephone Network - traditional telephone network
"""

GLOSSARY_SLIDE1 = GLOSSARY_COMMON + """- Direct Routing: Microsoft Teams feature connecting to PSTN via SBC
- Mediant VE: AudioCodes virtual SBC appliance for cloud/hypervisor deployment
"""

GLOSSARY_PROBLEM = GLOSSARY_COMMON + """- ec2:ReplaceRoute: AWS API action to modify a VPC route table entry
- ec2:AssociateAddress: AWS API action to assign an Elastic IP to an instance
- CloudTrail: AWS service recording API calls for auditing
- ARN: Amazon Resource Name - unique identifier for AWS resources
- Route Table: AWS networking construct directing traffic between subnets
"""

GLOSSARY_OPTION1 = GLOSSARY_COMMON + """- EventBridge: AWS serverless event bus for routing events to targets
- Lambda: AWS serverless compute service that runs code in response to events
- SSM Parameter Store: AWS service for storing configuration data and secrets
- CloudWatch: AWS monitoring and observability service
- SNS: Simple Notification Service - AWS managed messaging service
- Provisioned Concurrency: Lambda feature eliminating cold start latency
- AWS Config: Service for assessing, auditing, and evaluating AWS resource configurations
- VPC Flow Logs: Feature capturing IP traffic information for VPC network interfaces
- Stack Manager: AudioCodes management VM for deploying SBC HA stacks via CloudFormation
"""

GLOSSARY_OPTION2 = GLOSSARY_COMMON + """- CloudFormation: AWS infrastructure-as-code service for provisioning resources
- Stack Manager: AudioCodes management VM required for HA deployment (not standalone)
"""

GLOSSARY_OPTION3 = GLOSSARY_COMMON + """- DNS Failover: Using DNS record changes to redirect traffic to a backup system
- Configuration Drift: Divergence between configurations of systems meant to be identical
- Teams DR Priority/Weight: Microsoft Teams Direct Routing settings controlling SBC preference
"""

GLOSSARY_OPTION4 = GLOSSARY_COMMON + """- Mediant 800/1000/2600: AudioCodes physical SBC appliance models
- Rack-and-Stack: Physical installation of hardware in a data centre rack
"""

GLOSSARY_OPTION5 = GLOSSARY_COMMON + """- Mediant VE: AudioCodes Mediant Virtual Edition - software SBC for hypervisors
- Hypervisor: Software creating and managing virtual machines (VMware, Hyper-V, KVM)
- KVM: Kernel-based Virtual Machine - Linux virtualisation technology
"""

GLOSSARY_MATRIX = GLOSSARY_COMMON

GLOSSARY_RECOMMENDATION = GLOSSARY_COMMON + """- Configuration Drift: Divergence between configurations of systems meant to be identical
- Cloud-first: Organisational strategy prioritising cloud over on-premises deployment
"""

GLOSSARY_DECISION = GLOSSARY_COMMON + """- Cloud-first: Organisational strategy prioritising cloud over on-premises deployment
- Go-live: Point at which a system enters production use
"""


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, runs, alignment=PP_ALIGN.LEFT):
    """Add textbox with multiple styled runs. runs = [(text, size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, size, bold, color) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def add_bullet_textbox(slide, left, top, width, height, items, font_size=11,
                       color=DARK_TEXT, bold_items=None):
    """Add textbox with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    bold_items = bold_items or []
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.font.bold = i in bold_items
        p.space_after = Pt(4)
        # bullet
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)
        pPr.set('indent', str(Emu(Inches(0.25))))
        pPr.set('marL', str(Emu(Inches(0.35))))
    return txBox


def set_cell_fill(cell, color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def add_table(slide, left, top, width, height, rows_data, col_widths=None,
              header_bg=TABLE_HEADER_BG, font_size=9):
    """Add formatted table. rows_data = [[col1, col2, ...], ...]"""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0]) if rows_data else 0
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
                    p.font.bold = col_idx == 0

            if row_idx == 0:
                set_cell_fill(cell, (header_bg[0], header_bg[1], header_bg[2]))
            elif row_idx % 2 == 0:
                set_cell_fill(cell, (TABLE_ALT_BG[0], TABLE_ALT_BG[1], TABLE_ALT_BG[2]))
            else:
                set_cell_fill(cell, (TABLE_WHITE_BG[0], TABLE_WHITE_BG[1], TABLE_WHITE_BG[2]))

    return table_shape


def add_section_header(slide, text, top=Inches(0.3)):
    add_textbox(slide, Inches(0.5), top, Inches(9), Inches(0.5),
                text, font_size=22, bold=True, color=ACCENT_BLUE)


def add_subtitle_bar(slide, text, top):
    txBox = slide.shapes.add_textbox(Inches(0.5), top, Inches(9), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Calibri"
    return txBox


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Set author metadata
    prs.core_properties.author = "Kevin Sing"
    prs.core_properties.last_modified_by = "Kevin Sing"
    prs.core_properties.title = "AudioCodes SBC Deployment Options - Decision Pack"
    prs.core_properties.subject = "Microsoft Teams Direct Routing - SBC HA Architecture Options"
    prs.core_properties.created = datetime.datetime(2026, 3, 5)
    prs.core_properties.modified = datetime.datetime(2026, 3, 5)

    blank_layout = prs.slide_layouts[6]  # blank

    # =========================================================================
    # SLIDE 1 - Title
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, DARK_BG)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
                "AudioCodes SBC Deployment Options", font_size=36, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.6),
                "Decision Pack", font_size=28, bold=False, color=MID_GREY)
    add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
                "Microsoft Teams Direct Routing - SBC HA Architecture Options",
                font_size=16, color=MID_GREY)

    meta_rows = [
        ["Date", "5 March 2026"],
        ["Audience", "IT Manager, Cybersecurity, Cloud Platform, Voice Engineering"],
        ["Source", "AudioCodes SBC - Unified Deployment & Configuration Guide v2.6"],
        ["Purpose", "Present five SBC deployment options with recommendation"],
    ]
    y = Inches(4.5)
    for label, val in meta_rows:
        add_rich_textbox(slide, Inches(1), y, Inches(10), Inches(0.35), [
            (label + ":  ", 12, True, MID_GREY),
            (val, 12, False, WHITE),
        ])
        y += Inches(0.35)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_SLIDE1

    # =========================================================================
    # SLIDE 2 - The Problem
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "1. The Problem")

    add_textbox(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.8),
                "A cybersecurity review identified that ec2:ReplaceRoute - the AWS API used by AudioCodes SBCs "
                "for internal VIP failover - cannot be scoped to individual route entries. The SBC IAM policy is "
                "already at maximum IAM granularity (specific route table ARN + Env tag), but a compromised SBC "
                "could still replace any route in that table, not just VIP routes.",
                font_size=12, color=DARK_TEXT)

    add_subtitle_bar(slide, "Two Failover Paths - Only One Is Affected", Inches(1.9))

    table_data = [
        ["", "Internal (VIP)", "External (EIP)"],
        ["Failover API", "ec2:ReplaceRoute", "ec2:AssociateAddress"],
        ["What moves", "Route table entry to standby ENI", "EIP to standby WAN ENI"],
        ["Connects to", "Downstream SBCs, PBX, SIP providers", "Microsoft Teams"],
        ["IAM scoping", "Route table ARN + Env tag", "EIP ARN + App + Env tags"],
        ["Concern", "Cannot restrict to specific routes", "No concern - scoped to single EIP"],
    ]
    add_table(slide, Inches(0.5), Inches(2.3), Inches(8), Inches(2.5), table_data,
              col_widths=[Inches(1.8), Inches(3.1), Inches(3.1)], font_size=10)

    key_points = [
        "External EIP path (ec2:AssociateAddress) is not affected - already scoped to specific EIP with dual tag conditions",
        "Exploiting this requires the SBC instance to be compromised first - not an internet-facing attack surface",
        "Blast radius bounded to a single route table in a single tagged environment",
        "AWS has confirmed there are no condition keys for destinationCidrBlock or networkInterfaceId on ReplaceRoute",
    ]
    add_subtitle_bar(slide, "Key Points", Inches(5.0))
    add_bullet_textbox(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(1.8), key_points, font_size=10)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_PROBLEM

    # =========================================================================
    # SLIDE 3 - Option 1
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Option 1 - HA with Retrospective Guardrails")
    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(5), Inches(0.3),
                "RECOMMENDED", font_size=11, bold=True, color=ACCENT_GREEN)

    opt1_data = [
        ["Aspect", "Detail"],
        ["Architecture", "1+1 Active/Standby across two AZs - standard AudioCodes HA via Stack Manager"],
        ["Failover", "SBC firmware calls AWS APIs directly; VIP + EIP move to standby; active calls survive"],
        ["IAM permissions", "ec2:ReplaceRoute + ec2:AssociateAddress - both retained"],
        ["Security controls", "4-layer compensating control architecture (see right)"],
        ["Exposure window", "~7-18 seconds from malicious API call to full automated containment"],
        ["Guardrail cost", "~$6/month per region"],
        ["Licensing", "1x session/feature licence (single logical SBC) + 2x base VM licences"],
        ["Infrastructure", "2x SBC instances + 1x Stack Manager (t3.medium) per environment"],
    ]
    add_table(slide, Inches(0.3), Inches(1.1), Inches(6.5), Inches(3.5), opt1_data,
              col_widths=[Inches(1.8), Inches(4.7)], font_size=9)

    add_subtitle_bar(slide, "Compensating Controls (4 Layers)", Inches(1.1))
    # Position the subtitle on the right side
    slide.shapes[-1].left = Inches(7.0)
    slide.shapes[-1].width = Inches(6)

    controls = [
        "Dedicated VIP route table - only VIP routes in the table; blast radius limited to VIP entries only",
        "EventBridge + Lambda containment gate - validates every ReplaceRoute call; auto-reverts, strips IAM, quarantines, alerts",
        "Scheduled canary Lambda - polls route table every 60s against VIP allowlist; independent backstop",
        "AWS Config + VPC Flow Logs - custom Config rule flags non-VIP entries; Flow Log anomaly detection",
    ]
    add_bullet_textbox(slide, Inches(7.0), Inches(1.5), Inches(6), Inches(2.2), controls, font_size=9)

    extras = [
        "Why reactive, not preventive: AudioCodes firmware calls the EC2 API directly during failover - proprietary behaviour that cannot be intercepted",
        "Containment posture: Deny-by-default, fully automated, no human gate. Compromised instance left running in quarantine for forensics",
        "Cost breakdown: Containment Lambda ~$1, canary ~$1, provisioned concurrency ~$3, CloudWatch logs ~$1, EventBridge/SSM/SNS/SG = free tier",
    ]
    add_bullet_textbox(slide, Inches(7.0), Inches(3.8), Inches(6), Inches(1.8), extras, font_size=9)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_OPTION1

    # =========================================================================
    # SLIDE 4 - Option 2
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Option 2 - Standalone SBC (No HA)")
    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(5), Inches(0.3),
                "NOT RECOMMENDED", font_size=11, bold=True, color=ACCENT_RED)

    opt2_data = [
        ["Aspect", "Detail"],
        ["Architecture", "Single SBC instance per region - no HA pairing, no VIP, no route table manipulation"],
        ["Failover", "None - manual recovery (instance restart/replacement)"],
        ["IAM permissions", "ec2:ReplaceRoute removed entirely; ec2:AssociateAddress not required"],
        ["Security controls", "N/A - no IAM risk to mitigate"],
        ["Single point of failure", "Yes - SBC failure = total voice outage for region"],
        ["Licensing", "1x SBC licence per region"],
        ["Infrastructure", "1x SBC instance per region; no Stack Manager required"],
    ]
    add_table(slide, Inches(0.3), Inches(1.1), Inches(7), Inches(3.2), opt2_data,
              col_widths=[Inches(2), Inches(5)], font_size=10)

    add_subtitle_bar(slide, "Key Concerns", Inches(4.5))
    concerns = [
        "TOTAL REBUILD for future HA - most likely required at project completion. Kapila (AudioCodes) confirmed standalone to HA is a complete tear-down and rebuild",
        "Stack Manager must deploy the HA pair from scratch via CloudFormation - cannot retrofit HA onto existing standalone instance",
        "Recovery time on failure: minutes to hours. Could extend up to a day if the MSP needs to familiarise with a complex SBC recovery procedure",
        "All active calls drop on instance failure",
        "Fastest path to production - but paints the organisation into a corner",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(2.2), concerns,
                       font_size=10, bold_items=[0, 2])

    slide.notes_slide.notes_text_frame.text = GLOSSARY_OPTION2

    # =========================================================================
    # SLIDE 5 - Option 3
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Option 3 - Non-Seamless HA (2x Standalone SBCs)")

    opt3_data = [
        ["Aspect", "Detail"],
        ["Architecture", "Two independent SBC instances per region - not an AudioCodes HA pair"],
        ["Failover", "External: DNS-based, SIP proxy primary/secondary, or Teams DR priority/weight settings"],
        ["IAM permissions", "ec2:ReplaceRoute removed entirely; ec2:AssociateAddress removed entirely"],
        ["Security controls", "N/A - no IAM risk to mitigate"],
        ["Call survivability", "Active calls on primary drop - no session state synchronisation"],
        ["Licensing", "2x SBC licence per region (double cost)"],
        ["Infrastructure", "2x SBC instances per region; no Stack Manager required"],
    ]
    add_table(slide, Inches(0.3), Inches(1.0), Inches(7), Inches(3.2), opt3_data,
              col_widths=[Inches(2), Inches(5)], font_size=10)

    add_subtitle_bar(slide, "Key Considerations", Inches(4.4))
    notes3 = [
        "Eliminates all route table and EIP IAM permissions - noted IAM risk not present",
        "HA mechanism does not programmatically manipulate the AWS route table - no concerns with this approach",
        "Failover speed depends on DNS TTL or SIP retry timers - not seamless like true HA",
        "Each SBC maintains independent configuration - risk of configuration drift (ongoing operational burden)",
        "Not an AudioCodes-supported HA pattern - this is a custom resilience design",
        "Requires SIP provider and Microsoft coordination for primary/secondary configuration",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(4.8), Inches(12), Inches(2.3), notes3, font_size=10)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_OPTION3

    # =========================================================================
    # SLIDE 6 - Option 4
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Option 4 - On-Premises SBC in HA (Physical)")
    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(5), Inches(0.3),
                "ALTERNATIVE FALLBACK", font_size=11, bold=True, color=ACCENT_ORANGE)

    opt4_data = [
        ["Aspect", "Detail"],
        ["Architecture", "Physical AudioCodes Mediant appliance (800/1000/2600) with VRRP or native HA"],
        ["Failover", "Traditional HA - VRRP/AudioCodes native; active calls survive"],
        ["IAM permissions", "N/A - no cloud IAM involved"],
        ["Security controls", "Physical security, network ACLs - shifts to physical access domain"],
        ["Call survivability", "Active calls survive (session state synchronised)"],
        ["Licensing", "Hardware appliance licence"],
        ["Infrastructure", "Physical appliance + rack space + power + cooling"],
    ]
    add_table(slide, Inches(0.3), Inches(1.1), Inches(7), Inches(3.2), opt4_data,
              col_widths=[Inches(2), Inches(5)], font_size=10)

    add_subtitle_bar(slide, "Key Considerations", Inches(4.5))
    notes4 = [
        "Proven HA mechanism with no cloud API dependency - noted IAM risk not present",
        "Against cloud-first strategy - moves voice infrastructure back on-premises",
        "Hardware procurement lead time: weeks to months depending on model and availability",
        "Requires physical data centre presence with rack space, power, cooling, and network connectivity",
        "Ongoing hardware maintenance and lifecycle management (patching, warranty, end-of-life)",
        "Security concerns shift entirely from IAM/cloud to physical access controls and network segmentation",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(2.2), notes4, font_size=10)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_OPTION4

    # =========================================================================
    # SLIDE 7 - Option 5
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "Option 5 - On-Premises Virtualised SBC in HA")
    add_textbox(slide, Inches(0.5), Inches(0.75), Inches(5), Inches(0.3),
                "PREFERRED FALLBACK", font_size=11, bold=True, color=ACCENT_GREEN)

    opt5_data = [
        ["Aspect", "Detail"],
        ["Architecture", "AudioCodes Mediant VE (Virtual Edition) on existing on-premises hypervisors in HA"],
        ["Failover", "Traditional HA - VRRP/native AudioCodes HA on hypervisor; active calls survive"],
        ["IAM permissions", "N/A - no cloud IAM involved"],
        ["Security controls", "Existing on-premises security controls, hypervisor access controls, network ACLs"],
        ["Call survivability", "Active calls survive (session state synchronised)"],
        ["Licensing", "Mediant VE software licence (session/feature + base VM)"],
        ["Infrastructure", "2x VMs on existing hypervisors - no new hardware procurement"],
    ]
    add_table(slide, Inches(0.3), Inches(1.1), Inches(7), Inches(3.2), opt5_data,
              col_widths=[Inches(2), Inches(5)], font_size=10)

    add_subtitle_bar(slide, "Key Considerations", Inches(4.5))
    notes5 = [
        "Same Mediant VE software as the AWS deployment but hosted on existing on-premises hypervisors (VMware, Hyper-V, KVM)",
        "No hardware procurement lead time - uses existing hypervisor capacity, significantly faster than Option 4",
        "No new rack space, power, or cooling required - leverages existing data centre infrastructure",
        "Full AudioCodes HA support - same HA mechanism as Option 4 but virtualised",
        "Noted IAM risk not present - no cloud IAM permissions involved in failover",
        "Against cloud-first strategy - same strategic trade-off as Option 4, but lower capital cost and faster deployment",
        "Dependent on existing hypervisor capacity and availability in required data centre locations",
    ]
    add_bullet_textbox(slide, Inches(0.5), Inches(4.9), Inches(12), Inches(2.2), notes5, font_size=10)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_OPTION5

    # =========================================================================
    # SLIDE 8 - Comparison Matrix
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "3. Comparison Matrix")

    matrix_data = [
        ["Category", "Opt 1: HA + Guardrails", "Opt 2: Standalone", "Opt 3: 2x Standalone", "Opt 4: On-Prem Physical", "Opt 5: On-Prem Virtual"],
        ["HA capability", "Full 1+1 Active/Standby", "None", "External failover (DNS/SIP)", "Full HA (VRRP/native)", "Full HA (VRRP/native)"],
        ["Call survivability", "Active calls survive", "All calls drop", "Active calls drop", "Active calls survive", "Active calls survive"],
        ["Recovery time", "Seconds (automatic)", "Mins-hours, up to a day", "Secs-mins (DNS/SIP)", "Seconds (automatic)", "Seconds (automatic)"],
        ["Security risk", "7-18s window, 4-layer defence", "Noted IAM risk not present", "Noted IAM risk not present", "Noted IAM risk not present", "Noted IAM risk not present"],
        ["Build effort", "Moderate-high", "Low (fastest)", "Moderate", "Variable (hardware)", "Moderate - existing infra"],
        ["Licensing", "1x session/feature + 2x base VM + ~$6/mo", "1x SBC/region", "2x full SBC/region", "Appliance licence", "VE software licence"],
        ["Future flexibility", "Already at target state", "Total rebuild for HA", "Already HA - no concerns", "Locked to on-premises", "Locked to on-premises"],
        ["Vendor support", "AudioCodes-supported HA", "Standard support", "Custom pattern", "Hardware support", "AudioCodes-supported HA"],
        ["Cyber approval", "Depends on risk appetite", "High likelihood", "High likelihood", "High likelihood", "High likelihood"],
        ["Cloud-first aligned", "Yes", "Yes", "Yes", "No", "No"],
        ["Timeline to deploy", "Moderate", "Fastest", "Moderate", "Slowest", "Moderate-fast"],
    ]

    add_table(slide, Inches(0.2), Inches(0.9), Inches(12.9), Inches(5.5), matrix_data,
              col_widths=[Inches(1.7), Inches(2.24), Inches(2.24), Inches(2.24), Inches(2.24), Inches(2.24)],
              font_size=8)

    # Add key notes at bottom
    add_textbox(slide, Inches(0.3), Inches(6.5), Inches(12.5), Inches(0.8),
                "Options 1, 4 and 5 provide seamless failover with call survivability. Options 2 and 3 do not. "
                "Option 2 recovery could extend up to a day if the MSP needs to familiarise with a complex SBC recovery procedure. "
                "Choosing standalone (Option 2) now and needing HA later = total rebuild (AudioCodes confirmed).",
                font_size=9, color=DARK_TEXT)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_MATRIX

    # =========================================================================
    # SLIDE 9 - Recommendation
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "4. Recommendation")

    recs = [
        ("RECOMMENDED", ACCENT_GREEN, "Option 1 - HA with Retrospective Guardrails",
         "Delivers target-state HA architecture from day one. Avoids future rebuild risk. Seamless failover with "
         "call survivability. 7-18 second exposure window bounded by 4-layer automated defence. Guardrail cost ~$6/month per region."),

        ("PREFERRED FALLBACK", ACCENT_GREEN, "Option 5 - On-Premises Virtualised HA",
         "If Cyber rejects compensating controls and ec2:ReplaceRoute must be eliminated entirely. Full HA, call "
         "survivability, noted IAM risk not present. Uses Mediant VE on existing hypervisors - no hardware procurement. "
         "Trade-off: depends on cloud-first commitment."),

        ("ALTERNATIVE FALLBACK", ACCENT_ORANGE, "Option 4 - On-Premises Physical HA",
         "Same outcome as Option 5 using physical Mediant appliance. Use if hypervisor capacity not available. "
         "Slowest to deploy due to hardware procurement and data centre logistics."),

        ("ALSO CONSIDER", MID_GREY, "Option 3 - 2x Standalone SBCs",
         "Removes all IAM permissions. HA mechanism does not manipulate AWS route table. Active calls drop on primary "
         "failure. Configuration drift risk - each SBC maintains independent config. Ongoing operational burden."),

        ("NOT RECOMMENDED", ACCENT_RED, "Option 2 - Standalone (No HA)",
         "Single point of failure. Total rebuild required for future HA (most likely at project completion). Recovery "
         "could take up to a day. Paints the organisation into a corner."),
    ]

    y = Inches(1.0)
    for tag, tag_color, title, desc in recs:
        add_rich_textbox(slide, Inches(0.5), y, Inches(2.5), Inches(0.3), [
            (tag, 9, True, tag_color),
        ])
        add_rich_textbox(slide, Inches(2.5), y, Inches(10), Inches(0.3), [
            (title, 12, True, DARK_TEXT),
        ])
        add_textbox(slide, Inches(2.5), y + Inches(0.3), Inches(10), Inches(0.7),
                    desc, font_size=10, color=DARK_TEXT)
        y += Inches(1.15)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_RECOMMENDATION

    # =========================================================================
    # SLIDE 10 - Decision Framework & Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_section_header(slide, "5. Decision Framework and Next Steps")

    add_subtitle_bar(slide, "Factors to Score", Inches(0.85))
    factors_data = [
        ["#", "Factor", "Key Question"],
        ["1", "Risk appetite", "Is a 7-18s exposure window with automated containment acceptable, or must ReplaceRoute be eliminated entirely?"],
        ["2", "Voice criticality", "What is the acceptable RTO for voice services? Is a single point of failure tolerable?"],
        ["3", "Cloud-first commitment", "Does cloud-first strategy preclude an on-premises SBC?"],
        ["4", "Build vs. rebuild", "Is the risk of a total rebuild later (Option 2) acceptable given timeline pressures?"],
        ["5", "Cyber's position", "What is the formal position on compensating controls (Option 1)?"],
        ["6", "Timeline pressure", "How urgently must Direct Routing be in production?"],
        ["7", "Licensing budget", "Is double SBC licensing (Options 1 & 3) within budget?"],
    ]
    add_table(slide, Inches(0.3), Inches(1.2), Inches(12.5), Inches(2.8), factors_data,
              col_widths=[Inches(0.5), Inches(2), Inches(10)], font_size=9)

    add_subtitle_bar(slide, "Open Questions for the Decision Meeting", Inches(4.2))
    questions_data = [
        ["#", "Question", "Owner"],
        ["1", "What is the formal position on the 7-18s exposure window with 4-layer automated containment?", "Cybersecurity"],
        ["2", "Is the compensating control architecture sufficient to close the finding?", "Cybersecurity"],
        ["3", "If Option 1 is rejected, is the position permanent or reviewable after non-prod demonstration?", "Cybersecurity"],
        ["4", "What is the acceptable RTO for voice services?", "IT Management"],
        ["5", "Is the total rebuild risk (Option 2 to HA) acceptable?", "IT Management"],
        ["6", "Does cloud-first strategy rule out Options 4 and 5?", "IT Management"],
        ["7", "What is the estimated effort to build the guardrail stack (Option 1)?", "Cloud Platform / TTO"],
        ["8", "Can the guardrail stack be built and tested in non-prod before a final decision?", "Cloud Platform / TTO"],
        ["9", "Can regional SIP providers support primary/secondary SBC config (Option 3)?", "Voice Engineering"],
        ["10", "Is there sufficient hypervisor capacity on-premises for Option 5?", "Infra / Voice Eng"],
        ["11", "What is the hardware lead time for a physical AudioCodes appliance (Option 4)?", "Voice Engineering"],
        ["12", "What is the target date for Microsoft Teams Direct Routing go-live?", "All stakeholders"],
        ["13", "Is this decision per-region or must it be consistent across all regions?", "All stakeholders"],
    ]
    add_table(slide, Inches(0.3), Inches(4.55), Inches(12.5), Inches(2.8), questions_data,
              col_widths=[Inches(0.5), Inches(9.5), Inches(2.5)], font_size=8)

    slide.notes_slide.notes_text_frame.text = GLOSSARY_DECISION

    # Save
    output_path = "/home/kevin/projects/public/analyst-reports/SBC-Deployment-Options-Decision-Pack.pptx"
    prs.save(output_path)
    print(f"Saved to {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
